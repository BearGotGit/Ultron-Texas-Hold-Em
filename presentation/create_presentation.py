"""
Create PowerPoint presentation for Texas Hold'em Ultron project.
This script generates a comprehensive presentation covering the RL poker bot.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        # Add subtitle
        sub_box = slide.shapes.add_textbox(left, Inches(4), width, Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 102, 204)
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items, code_snippet=None):
    """Add a content slide with bullet points and optional code."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Content area
    if code_snippet:
        content_width = Inches(6)
        code_left = Inches(6.5)
    else:
        content_width = Inches(12)
        code_left = None
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle indentation levels
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = f"• {item}"
            p.level = 0
        
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Add code snippet if provided
    if code_snippet:
        code_box = slide.shapes.add_textbox(code_left, Inches(1.3), Inches(6.3), Inches(5.5))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_snippet
        p.font.size = Pt(11)
        p.font.name = "Courier New"
        
        # Add background to code
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide

def add_code_slide(title, code_snippet, description=""):
    """Add a slide focused on code."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.italic = True
    
    # Code box
    top = Inches(1.7) if description else Inches(1.2)
    code_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(5.3))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_snippet
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
    p.font.color.rgb = RGBColor(171, 178, 191)
    
    return slide

def add_diagram_slide(title, content_text):
    """Add a slide for diagrams/architecture."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Content (ASCII diagram or description)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content_text
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    return slide


# ==========================================
# SLIDE 1: Title Slide
# ==========================================
add_title_slide(
    "🃏 Texas Hold'em Ultron",
    "A Reinforcement Learning Poker Bot\n\nAnthony • Berend • Daniel • Dina • Eby • Aaron\nCS 4444 Project"
)

# ==========================================
# SECTION: Problem Definition
# ==========================================
add_section_slide("1. Problem Definition")

add_content_slide(
    "Project Goal",
    [
        "Build an AI agent that can play competitive Texas Hold'em poker",
        "Understand game state and board dynamics",
        "Evaluate hand strength and equity in real-time",
        "Predict opponent behavior and ranges",
        "Make optimal decisions: fold, call, or raise",
        "Compete against other teams' bots in class tournament"
    ]
)

add_content_slide(
    "Project Requirements",
    [
        "Real-time decision making (< 5 seconds per action)",
        "Support for 2-9 players per table",
        "Handle all poker stages: pre-flop, flop, turn, river",
        "Proper betting mechanics: blinds, raises, all-in situations",
        "Trained model that improves over time",
        "Evaluation metrics for performance tracking"
    ]
)

# ==========================================
# SECTION: Game Background
# ==========================================
add_section_slide("2. Game Background")

add_content_slide(
    "Texas Hold'em Poker Hands",
    [
        "🏆 Royal Flush - A, K, Q, J, 10 same suit (best hand)",
        "⭐ Straight Flush - Five sequential cards, same suit",
        "🎰 Four of a Kind - Four cards of same rank",
        "🏠 Full House - Three of a kind + a pair",
        "♠️ Flush - Five cards of same suit",
        "📏 Straight - Five sequential cards",
        "🎲 Three of a Kind - Three cards of same rank",
        "👥 Two Pair - Two different pairs",
        "👤 Pair - Two cards of same rank",
        "📝 High Card - Highest card wins (worst hand)"
    ]
)

add_content_slide(
    "Challenges in Poker AI",
    [
        "Imperfect Information - Can't see opponent's cards",
        "Probabilistic Outcomes - Must handle uncertainty",
        "Large State Space - 52 cards, multiple stages, varying pot sizes",
        "Opponent Modeling - Different players have different strategies",
        "Bluffing - Sometimes optimal play involves deception",
        "Multi-street Planning - Decisions affect future betting rounds",
        "Pot Odds vs Equity - Balancing risk and reward"
    ]
)

# ==========================================
# SECTION: System Design
# ==========================================
add_section_slide("3. System Design")

add_content_slide(
    "Why Reinforcement Learning?",
    [
        "Self-play enables learning without human examples",
        "Handles sequential decision making naturally",
        "Can adapt to different opponent styles",
        "PPO provides stable training for complex action spaces",
        "Continuous learning - model improves with more games",
        "No need for labeled poker hand databases"
    ]
)

add_diagram_slide(
    "Component Diagram",
    """
┌─────────────────────────────────────────────────────────────────────────────┐
│                           TEXAS HOLD'EM ULTRON                               │
├─────────────────────────────────────────────────────────────────────────────┤
│                                                                              │
│  ┌──────────────┐    ┌──────────────┐    ┌──────────────┐                   │
│  │   Agents     │    │  Simulation  │    │   Training   │                   │
│  ├──────────────┤    ├──────────────┤    ├──────────────┤                   │
│  │ PokerPlayer  │◄──►│  PokerEnv    │◄──►│ PPOTrainer   │                   │
│  │ MonteCarloAgt│    │  (Gymnasium) │    │              │                   │
│  │ RLAgent      │    │              │    │ RolloutBuffer│                   │
│  │ HumanPlayer  │    │ CardUtils    │    │              │                   │
│  └──────────────┘    └──────────────┘    │ PPOConfig    │                   │
│                             │            └──────────────┘                   │
│                             │                   │                           │
│                             ▼                   ▼                           │
│                      ┌──────────────┐    ┌──────────────┐                   │
│                      │  PokerPPO    │    │    Utils     │                   │
│                      │   Model      │    ├──────────────┤                   │
│                      ├──────────────┤    │ TensorBoard  │                   │
│                      │ CardEmbedding│    │ Reader       │                   │
│                      │ HandEmbedding│    │              │                   │
│                      │ FoldHead     │    │ Device Util  │                   │
│                      │ BetHead      │    └──────────────┘                   │
│                      │ ValueHead    │                                       │
│                      └──────────────┘                                       │
│                                                                              │
└─────────────────────────────────────────────────────────────────────────────┘
"""
)

add_diagram_slide(
    "Action Space Design",
    """
┌─────────────────────────────────────────────────────────────────────────────┐
│                           ACTION SPACE (2D Continuous)                       │
├─────────────────────────────────────────────────────────────────────────────┤
│                                                                              │
│  Action = [p_fold, bet_scalar]  where both ∈ [0, 1]                         │
│                                                                              │
│  ┌───────────────────────────────────────────────────────────────────┐      │
│  │  p_fold (Bernoulli Distribution)                                   │      │
│  │  ┌─────────────────────────────────────────────────────────────┐  │      │
│  │  │  > 0.5  →  FOLD (if there's a bet to call)                  │  │      │
│  │  │  ≤ 0.5  →  Continue to bet_scalar interpretation            │  │      │
│  │  └─────────────────────────────────────────────────────────────┘  │      │
│  └───────────────────────────────────────────────────────────────────┘      │
│                                                                              │
│  ┌───────────────────────────────────────────────────────────────────┐      │
│  │  bet_scalar (Beta Distribution)                                    │      │
│  │  ┌─────────────────────────────────────────────────────────────┐  │      │
│  │  │  < 0.1 (ε)      →  CHECK (if no bet) or CALL                │  │      │
│  │  │  0.1 - 0.9      →  RAISE (scaled between min_raise & stack) │  │      │
│  │  │  > 0.9 (1-ε)    →  ALL-IN                                   │  │      │
│  │  └─────────────────────────────────────────────────────────────┘  │      │
│  └───────────────────────────────────────────────────────────────────┘      │
│                                                                              │
│  Why Beta Distribution?                                                      │
│  • Naturally bounded [0, 1] - perfect for bet sizing                        │
│  • Flexible shape (α, β parameters) - can be uniform, skewed, peaked        │
│  • Differentiable - enables gradient-based learning                         │
│                                                                              │
└─────────────────────────────────────────────────────────────────────────────┘
"""
)

add_content_slide(
    "Core Design Decisions",
    [
        "Gymnasium Interface - Standard RL environment API",
        "Separate Embedding Modules - Cards, hand features, numeric features",
        "Hybrid Action Space - Bernoulli (fold) + Beta (bet sizing)",
        "Log-normalized Observations - Prevents gradient explosion",
        "Self-play Training - Hero vs Monte Carlo opponents",
        "TensorBoard Integration - Real-time training monitoring"
    ]
)

# ==========================================
# SECTION: Technical Implementation
# ==========================================
add_section_slide("4. Technical Implementation")

add_diagram_slide(
    "Project File Architecture",
    """
Ultron-Texas-Hold-Em/
│
├── agents/                          # Player implementations
│   ├── poker_player.py              # Abstract base class (PokerPlayer)
│   ├── monte_carlo_agent.py         # MC equity-based agent (opponent)
│   ├── rl_agent.py                  # RL agent wrapper
│   ├── human_player.py              # Interactive human player
│   └── agent.py                     # Legacy agent with equity calc
│
├── simulation/                      # Game environment
│   ├── poker_env.py                 # Gymnasium RL environment
│   ├── poker_simulator.py           # Game flow simulation
│   └── card_utils.py                # Card encoding utilities
│
├── training/                        # RL training code
│   ├── ppo_model.py                 # Neural network architecture
│   ├── train_rl_model.py            # PPO training loop
│   ├── losses.py                    # Custom loss functions
│   └── evaluate_model.py            # Model evaluation
│
├── utils/                           # Utilities
│   ├── tensorboard_reader.py        # Training log analysis
│   └── device.py                    # GPU/CPU detection
│
├── tests/                           # Test suite
│   ├── test_poker_env.py            # Environment tests
│   ├── test_model.py                # Model tests
│   └── ...
│
├── main.py                          # Interactive poker game
└── play_vs_rl.py                    # Play against trained model
"""
)

add_code_slide(
    "RL Foundation: Poker Environment (poker_env.py)",
    """class PokerEnv(gym.Env):
    \"\"\"Gymnasium-compatible Texas Hold'em environment.\"\"\"
    
    # Observation Space: 423 dimensions
    # - 7 x 53 card one-hot encodings (hole cards + board) = 371
    # - 10 binary hand features (pair, flush, etc.)        = 10
    # - 9 x 4 player features (money, bet, folded, all_in) = 36
    # - 6 global features (pot, current_bet, etc.)         = 6
    
    observation_space = spaces.Box(low=-inf, high=inf, shape=(423,))
    
    # Action Space: [p_fold, bet_scalar] both in [0, 1]
    action_space = spaces.Box(low=0.0, high=1.0, shape=(2,))
    
    def step(self, action):
        p_fold, bet_scalar = action[0], action[1]
        
        # Interpret action into poker move
        poker_action = interpret_action(
            p_fold=p_fold,
            bet_scalar=bet_scalar,
            current_bet=self.current_bet,
            my_bet=hero.bet,
            min_raise=self.min_raise,
            my_money=hero.money,
        )
        
        # Apply action, advance game
        self._apply_action(self.hero_idx, poker_action)
        self._advance_to_hero_or_end()
        
        reward = self._calculate_reward()  # Normalized chip delta
        return obs, reward, terminated, truncated, info""",
    "Gymnasium-compatible environment with continuous action space"
)

add_code_slide(
    "Deep Learning Model: PPO Architecture (ppo_model.py)",
    """class PokerPPOModel(nn.Module):
    \"\"\"PPO Actor-Critic for Texas Hold'em.\"\"\"
    
    def __init__(self, card_embed_dim=64, hidden_dim=256):
        # Card Embedding: 53-dim one-hot → 64-dim (shared across 7 cards)
        self.card_embedding = CardEmbedding(card_embed_dim)  # 53 → 64 → 64
        
        # Hand Embedding: 10 binary flags → 32-dim
        self.hand_embedding = nn.Sequential(
            nn.Linear(10, 32), nn.ReLU(), nn.Linear(32, 32))
        
        # Numeric Embedding: 42 features → 64-dim
        self.numeric_embedding = nn.Sequential(
            nn.Linear(42, 64), nn.ReLU(), nn.Linear(64, 64))
        
        # Combined: 448 + 32 + 64 = 544 → Shared Trunk → 256
        self.shared_trunk = nn.Sequential(
            nn.Linear(544, 256), nn.ReLU(),
            nn.Linear(256, 256), nn.ReLU())
        
        # Three Output Heads:
        self.fold_head = nn.Sequential(...)   # → 1 (Bernoulli logit)
        self.bet_head = nn.Sequential(...)    # → 2 (Beta α, β)
        self.value_head = nn.Sequential(...)  # → 1 (V(s) critic)
    
    def forward(self, obs):
        # Embed each feature type separately
        card_features = self._embed_cards(obs)    # (batch, 448)
        hand_features = self.hand_embedding(...)  # (batch, 32)
        numeric_features = self.numeric_embedding(...)  # (batch, 64)
        
        combined = torch.cat([card_features, hand_features, numeric_features], dim=1)
        hidden = self.shared_trunk(combined)
        
        fold_logit = self.fold_head(hidden)
        bet_alpha = F.softplus(self.bet_head(hidden)[:, 0:1]) + 1.0
        bet_beta = F.softplus(self.bet_head(hidden)[:, 1:2]) + 1.0
        value = self.value_head(hidden)
        
        return fold_logit, bet_alpha, bet_beta, value""",
    "Neural network with separate embeddings and actor-critic heads"
)

add_code_slide(
    "Monte Carlo Equity Component (monte_carlo_agent.py)",
    """class MonteCarloAgent(PokerPlayer):
    \"\"\"Uses Monte Carlo simulation for equity estimation.\"\"\"
    
    def _calculate_equity(self, hole_cards, board, num_opponents):
        \"\"\"Monte Carlo equity calculation.\"\"\"
        remaining = [c for c in full_deck if c not in used_cards]
        cards_to_deal = 5 - len(board)
        
        wins, ties, total = 0, 0, 0
        
        for _ in range(self.num_simulations):  # Default: 500 simulations
            # Sample random opponent hands and remaining board cards
            sampled = random.sample(remaining, 2*num_opponents + cards_to_deal)
            opponent_hands = [sampled[i*2:(i+1)*2] for i in range(num_opponents)]
            full_board = board + sampled[2*num_opponents:]
            
            # Evaluate all hands using Treys library
            my_score = self.evaluator.evaluate(full_board, hole_cards)
            opponent_scores = [self.evaluator.evaluate(full_board, h) 
                              for h in opponent_hands]
            
            best_opponent = min(opponent_scores)  # Lower is better in Treys
            
            if my_score < best_opponent:
                wins += 1
            elif my_score == best_opponent:
                ties += 0.5
            total += 1
        
        return (wins + ties) / total  # Equity: probability of winning
    
    def _make_decision(self, equity, pot_odds, to_call, pot, min_raise, my_money):
        \"\"\"Decision based on equity vs pot odds.\"\"\"
        if equity >= raise_threshold:
            return PokerAction.raise_to(...)
        elif equity > pot_odds:  # +EV to call
            return PokerAction.call(...)
        else:
            return PokerAction.fold()""",
    "Monte Carlo simulation for hand strength estimation"
)

add_code_slide(
    "PPO Training Method (train_rl_model.py)",
    """class PPOTrainer:
    def _update_policy(self):
        \"\"\"PPO policy update with clipped objective.\"\"\"
        for epoch in range(self.config.num_epochs):  # 4 epochs
            for batch in self.buffer.get_minibatches(self.config.num_minibatches):
                mb_obs, mb_actions, mb_log_probs, mb_advantages, mb_returns, mb_values = batch
                
                # Get new log probs and values
                new_log_probs, entropy, new_values = self.model.evaluate_actions(
                    mb_obs, mb_actions)
                
                # PPO ratio: π_new(a|s) / π_old(a|s)
                ratio = torch.exp(new_log_probs - mb_log_probs)
                
                # Clipped surrogate objective
                pg_loss1 = -mb_advantages * ratio
                pg_loss2 = -mb_advantages * torch.clamp(
                    ratio, 1 - self.config.clip_coef, 1 + self.config.clip_coef)
                pg_loss = torch.max(pg_loss1, pg_loss2).mean()
                
                # Value loss (with optional clipping)
                v_loss = 0.5 * ((new_values - mb_returns) ** 2).mean()
                
                # Entropy bonus for exploration
                entropy_loss = entropy.mean()
                
                # Total loss: policy - entropy_bonus + value
                loss = (pg_loss - self.config.ent_coef * entropy_loss
                        + self.config.vf_coef * v_loss)
                
                # Optimize
                self.optimizer.zero_grad()
                loss.backward()
                nn.utils.clip_grad_norm_(self.model.parameters(), 0.5)
                self.optimizer.step()""",
    "PPO with clipped objective, value function, and entropy regularization"
)

add_code_slide(
    "Training Pipeline Overview",
    """# train_rl_model.py - Main Training Loop

config = PPOConfig(
    total_timesteps=1_000_000,  # Total training decisions
    learning_rate=3e-4,
    gamma=0.99,           # Discount factor
    gae_lambda=0.95,      # GAE lambda
    num_steps=128,        # Steps per rollout
    num_epochs=4,         # PPO epochs per update
    clip_coef=0.2,        # PPO clipping
    ent_coef=0.01,        # Entropy coefficient
    vf_coef=0.5,          # Value function coefficient
)

trainer = PPOTrainer(config)

# Main loop
for iteration in range(num_iterations):
    # 1. Collect rollouts (hero vs MC opponents)
    rollout_stats = trainer._collect_rollouts()
    
    # 2. Compute GAE advantages
    buffer.compute_returns_and_advantages(last_value, gamma, gae_lambda)
    
    # 3. Update policy with PPO
    update_stats = trainer._update_policy()
    
    # 4. Log to TensorBoard
    writer.add_scalar("rollout/mean_reward", rollout_stats["mean_reward"])
    writer.add_scalar("losses/pg_loss", update_stats["pg_loss"])
    
    # 5. Periodic evaluation against fixed opponents
    if iteration % eval_interval == 0:
        eval_stats = trainer.evaluate(num_episodes=100)
        print(f"Win rate: {eval_stats['win_rate']:.2%}")
    
    # 6. Save checkpoints
    if iteration % save_interval == 0:
        trainer.save_checkpoint(f"checkpoint_{iteration}.pt")""",
    "Complete training pipeline with rollout collection, PPO updates, and evaluation"
)

add_code_slide(
    "Observation Encoding (poker_env.py)",
    """# Observation normalization for stable training

def _get_observation(self) -> np.ndarray:
    obs_parts = []
    
    # 1. Card encodings (7 x 53 one-hot)
    for card in hole_cards + board:
        obs_parts.append(encode_card_one_hot(card))  # 53-dim per card
    
    # 2. Hand features (10 binary flags)
    hand_features = encode_hand_features(hole_cards, board)
    obs_parts.append(hand_features)  # pair, two_pair, flush, etc.
    
    # 3. Player features (log-normalized)
    for player in self.players:
        stack_normalizer = np.log1p(self.config.starting_stack)
        features = [
            np.log1p(player.money) / stack_normalizer,  # Log-normalized stack
            np.log1p(player.bet) / np.log1p(big_blind), # Log-normalized bet
            float(player.folded),   # Binary
            float(player.all_in),   # Binary
        ]
    
    # 4. Global features (normalized)
    global_features = [
        np.log1p(self.pot.money) / np.log1p(total_starting_money),
        np.log1p(self.current_bet) / np.log1p(big_blind),
        np.log1p(self.min_raise) / np.log1p(big_blind),
        encode_round_stage(self.round_stage) / 4.0,  # 0-1 normalized
        self.hero_idx / max(num_players - 1, 1),     # Position
        self.dealer_position / max(num_players - 1, 1),
    ]
    
    return np.concatenate(obs_parts)  # Total: 423 dimensions""",
    "Log-normalization prevents gradient explosion from large chip values"
)

add_code_slide(
    "Action Interpretation (poker_env.py)",
    """def interpret_action(p_fold, bet_scalar, current_bet, my_bet, 
                         min_raise, my_money) -> PokerAction:
    \"\"\"Convert neural network output to poker action.\"\"\"
    
    to_call = current_bet - my_bet
    
    # Fold decision (Bernoulli)
    if p_fold > 0.5:
        if to_call > 0:
            return PokerAction.fold()
        else:
            # Can't fold when nothing to call - CHECK instead
            return PokerAction.check()
    
    # Bet sizing (Beta distribution output)
    ACTION_EPSILON = 0.1
    
    if bet_scalar < ACTION_EPSILON:
        # Low bet_scalar → Check or Call
        if to_call <= 0:
            return PokerAction.check()
        return PokerAction.call(min(to_call, my_money))
    
    elif bet_scalar > 1.0 - ACTION_EPSILON:
        # High bet_scalar → All-in
        return PokerAction.raise_to(my_money)
    
    else:
        # Middle values → Scaled raise
        normalized = (bet_scalar - ACTION_EPSILON) / (1.0 - 2*ACTION_EPSILON)
        raise_amount = int(min_raise + normalized * (my_money - min_raise))
        raise_amount = max(min_raise, min(raise_amount, my_money))
        return PokerAction.raise_to(raise_amount)""",
    "Smooth mapping from [0,1] to poker actions"
)

# ==========================================
# SECTION: Experiments and Evaluation
# ==========================================
add_section_slide("5. Experiments & Evaluation")

add_content_slide(
    "Evaluation Methodology",
    [
        "Self-play against Monte Carlo opponents with varying aggression",
        "Periodic evaluation every 50 training iterations",
        "Metrics tracked via TensorBoard:",
        ("• Win rate against baseline opponents", 1),
        ("• Average profit per hand", 1),
        ("• Episode length (actions per hand)", 1),
        ("• Policy entropy (exploration measure)", 1),
        "Checkpoint saving for model versioning",
        "Debug mode for step-by-step action analysis"
    ]
)

add_content_slide(
    "Key Metrics",
    [
        "rollout/mean_reward - Average reward per episode",
        "rollout/mean_length - Actions per hand (should be > 1)",
        "eval/win_rate - Win percentage against opponents",
        "eval/avg_profit - Average chip gain/loss per hand",
        "losses/pg_loss - Policy gradient loss",
        "losses/value_loss - Value function MSE",
        "losses/entropy - Action distribution entropy",
        "losses/kl_div - KL divergence (policy stability)"
    ]
)

add_code_slide(
    "Evaluation Code (train_rl_model.py)",
    """def evaluate(self, num_episodes: int) -> Dict[str, float]:
    \"\"\"Evaluate the current policy against opponents.\"\"\"
    self.model.eval()
    
    wins = 0
    total_profit = 0.0
    env = self.envs[0]  # Use first environment
    
    for _ in range(num_episodes):
        obs, _ = env.reset()
        initial_chips = env.players[env.hero_idx].money
        done = False
        
        while not done:
            obs_t = torch.tensor(obs, dtype=torch.float32).unsqueeze(0)
            with torch.no_grad():
                action, _, _, _ = self.model.get_action_and_value(
                    obs_t, deterministic=True)  # Use mean, not sample
            
            obs, reward, terminated, truncated, info = env.step(action)
            done = terminated or truncated
        
        final_chips = env.players[env.hero_idx].money
        profit = final_chips - initial_chips
        
        if profit > 0:
            wins += 1
        total_profit += profit
    
    return {
        "win_rate": wins / num_episodes,
        "avg_profit": total_profit / num_episodes,
    }""",
    "Deterministic evaluation against Monte Carlo baseline"
)

add_content_slide(
    "Training Results & Observations",
    [
        "Architecture improvements (normalized inputs) prevent always-fold",
        "Episode length > 1.0 indicates meaningful play",
        "Win rate baseline: ~50% against random, target: 55%+",
        "Key insight: Log-normalization critical for monetary values",
        "Separate embeddings help model learn feature interactions",
        "Beta distribution provides smooth bet sizing"
    ]
)

# ==========================================
# SECTION: Innovation Highlights
# ==========================================
add_section_slide("6. Innovation Highlights")

add_content_slide(
    "Innovation Highlights",
    [
        "🎯 Hybrid Action Space - Bernoulli fold + Beta bet sizing",
        "🔧 Three-Head Architecture - Separate card/hand/numeric embeddings",
        "📊 Log-Normalized Observations - Prevents gradient explosion",
        "🎮 Gymnasium Interface - Standard RL environment compatibility",
        "🤖 Monte Carlo Bootstrapping - Quality opponents from day 1",
        "📈 TensorBoard Integration - Real-time training visualization",
        "🔄 Self-Play Pipeline - Continuous improvement loop",
        "🎲 Flexible Player Support - 2-9 players per table"
    ]
)

add_content_slide(
    "Technical Innovations",
    [
        "Card Embedding Module - Shared weights across all 7 card slots",
        "Hand Feature Extraction - Binary flags for made hands (pair, flush...)",
        "GAE (Generalized Advantage Estimation) - Reduced variance",
        "Orthogonal Initialization - Stable early training",
        "Small Gain for Output Heads - Prevents p_fold saturation",
        "Clipped Value Loss - More stable critic updates"
    ]
)

add_content_slide(
    "Future Improvements",
    [
        "Opponent modeling - Adapt to different play styles",
        "CFR integration - Counterfactual regret for Nash equilibrium",
        "Deeper networks - Transformer architecture for card sequences",
        "League training - Population of diverse opponents",
        "Multi-table support - Parallel environment scaling",
        "Real money API integration - Deploy to online platforms"
    ]
)

# ==========================================
# FINAL SLIDE
# ==========================================
add_title_slide(
    "Thank You! 🃏",
    "Questions?\n\nGitHub: github.com/BearGotGit/Ultron-Texas-Hold-Em"
)

# Save presentation
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Texas_Holdem_Ultron_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
